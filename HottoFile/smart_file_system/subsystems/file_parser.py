import os
import subprocess
from datetime import datetime

import pytesseract
from PIL import Image

import pdfplumber
import docx
import pptx
import openpyxl

from data_structures.file_info import FileInfo


class FileParser:
    def __init__(self, llm_model="minicpm-v4.5"):
        """
        :param llm_model: Ollama 本地运行的模型名称
        """
        self.llm_model = llm_model

    def _call_ollama(self, prompt, image_path=None):
        """
        调用 Ollama 的本地大模型 (minicpm-v4.5)，支持文本或图片输入
        """
        try:
            cmd = ["ollama", "run", self.llm_model, prompt]
            if image_path:
                cmd.extend(["--image", image_path])
            result = subprocess.run(cmd, capture_output=True, text=True)
            return result.stdout.strip()
        except Exception as e:
            print(f"调用 Ollama 模型失败: {e}")
            return ""

    def parse_file(self, file_path):
        try:
            file_name = os.path.basename(file_path)
            file_ext = file_name.split('.')[-1].lower()
            file_size = os.path.getsize(file_path)
            created_at = datetime.fromtimestamp(os.path.getctime(file_path)).isoformat()
            modified_at = datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat()

            text_summary = None
            image_features = None
            ftype = "other"

            # 📄 文本类文件
            if file_ext in ["txt", "md"]:
                ftype = "document"
                with open(file_path, "r", encoding="utf-8") as f:
                    raw_text = f.read()
                # 交给本地模型总结
                text_summary = self._call_ollama(f"请总结以下文本内容：\n{raw_text[:2000]}")

            elif file_ext == "pdf":
                ftype = "document"
                with pdfplumber.open(file_path) as pdf:
                    raw_text = "\n".join([page.extract_text() or "" for page in pdf.pages])
                text_summary = self._call_ollama(f"请总结以下 PDF 内容：\n{raw_text[:2000]}")

            elif file_ext in ["docx", "doc"]:
                ftype = "document"
                doc = docx.Document(file_path)
                raw_text = "\n".join([p.text for p in doc.paragraphs])
                text_summary = self._call_ollama(f"请总结以下 Word 文档内容：\n{raw_text[:2000]}")

            elif file_ext in ["pptx", "ppt"]:
                ftype = "document"
                prs = pptx.Presentation(file_path)
                raw_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            raw_text.append(shape.text)
                raw_text = "\n".join(raw_text)
                text_summary = self._call_ollama(f"请总结以下 PPT 内容：\n{raw_text[:2000]}")

            elif file_ext in ["xlsx", "xls"]:
                ftype = "document"
                wb = openpyxl.load_workbook(file_path)
                raw_text = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        raw_text.append(" ".join([str(cell) if cell else "" for cell in row]))
                raw_text = "\n".join(raw_text)
                text_summary = self._call_ollama(f"请总结以下 Excel 表格内容：\n{raw_text[:2000]}")

            # 🖼️ 图片文件
            elif file_ext in ["jpg", "jpeg", "png"]:
                ftype = "image"
                try:
                    # OCR 提取文本
                    ocr_text = pytesseract.image_to_string(Image.open(file_path))
                    # 用视觉模型生成图像描述
                    image_desc = self._call_ollama("请描述这张图片的内容。", image_path=file_path)
                    image_features = image_desc
                    text_summary = f"OCR文本：{ocr_text}\n视觉模型描述：{image_desc}"
                except Exception as e:
                    print(f"OCR/图像解析失败: {e}")
                    text_summary = ""
                    image_features = ""

            # 🎵 音频 & 🎬 视频 未来可加

            # 组装 FileInfo
            return FileInfo(
                path=file_path,
                name=file_name,
                ext=file_ext,
                ftype=ftype,
                size=file_size,
                created_at=created_at,
                modified_at=modified_at,
                text_summary=text_summary,
                image_features=image_features
            )

        except Exception as e:
            print(f"解析文件 {file_path} 失败: {e}")
            return None
